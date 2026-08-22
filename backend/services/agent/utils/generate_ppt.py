import io
from pptx import Presentation


async def generate_ppt(data: dict) -> bytes:
    """
    Generates a PowerPoint presentation (.pptx) from JSON data.
    Expected data structure:
    {
        "title": "Main Title",
        "subtitle": "Main Subtitle",
        "slides": [
            {
                "title": "Slide Title",
                "points": ["Point 1", "Point 2", ...]
            },
            ...
        ]
    }
    """
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = data.get("title", "Presentation")
    if data.get("subtitle"):
        subtitle.text = data.get("subtitle", "")

    # Bullet Points Slides
    bullet_slide_layout = prs.slide_layouts[1]
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = slide_data.get("title", "")

        tf = body_shape.text_frame
        points = slide_data.get("points", [])
        if points:
            tf.text = points[0]
            for point in points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

    # Save presentation to BytesIO
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer.getvalue()